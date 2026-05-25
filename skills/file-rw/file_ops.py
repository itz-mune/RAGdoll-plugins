"""
Filesystem operations for the File R/W skill.

Every write / destructive operation is designed to be called in two phases:
  1. preview()  — computes diff / stats, returns data for the permission popup
  2. execute()  — called after the user approves; writes atomically

All public methods are async so they can be awaited from _arun without
running heavy I/O on the event loop thread (uses run_in_executor internally).
"""
from __future__ import annotations

import asyncio
import mimetypes
import os
import shutil
import stat
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional


# ── Result types ───────────────────────────────────────────────────────────────

@dataclass
class FileReadResult:
    content:      str
    encoding:     str
    size_bytes:   int
    line_count:   int
    was_truncated: bool
    truncated_at:  int = 0   # byte offset where truncation happened


@dataclass
class WriteResult:
    ok:      bool
    path:    str
    error:   Optional[str] = None
    backup_path: Optional[str] = None


@dataclass
class OperationResult:
    ok:    bool
    path:  str
    error: Optional[str] = None
    extra: dict = field(default_factory=dict)


@dataclass
class FileStats:
    name:             str
    absolute_path:    str
    size_bytes:       int
    size_human:       str
    is_file:          bool
    is_directory:     bool
    is_symlink:       bool
    is_hidden:        bool
    extension:        str
    mime_type:        str
    created_at:       str    # ISO 8601
    modified_at:      str
    accessed_at:      str
    permissions_octal: str   # "0644"
    permissions_human: str   # "rw-r--r--"
    owner:             str
    is_readable:       bool
    is_writable:       bool
    is_executable:     bool
    child_count:       Optional[int] = None   # dirs only
    total_size:        Optional[int] = None   # dirs only (recursive)

    def to_dict(self) -> dict:
        return {k: v for k, v in self.__dict__.items()}


@dataclass
class DirectoryEntry:
    name:         str
    is_directory: bool
    size_bytes:   Optional[int]
    modified_at:  str
    extension:    str
    is_hidden:    bool


@dataclass
class DirectoryListing:
    path:         str
    entries:      list[DirectoryEntry]
    total_count:  int
    hidden_count: int


# ── Encoding detection ──────────────────────────────────────────────────────────

def _decode_bytes(raw: bytes) -> tuple[str, str]:
    """Try UTF-8, then chardet, then latin-1 fallback. Returns (text, encoding)."""
    try:
        return raw.decode("utf-8"), "utf-8"
    except UnicodeDecodeError:
        pass
    try:
        import chardet
        detected = chardet.detect(raw[:32_768])  # sample for speed
        enc = detected.get("encoding") or "latin-1"
        return raw.decode(enc, errors="replace"), enc
    except ImportError:
        pass
    return raw.decode("latin-1", errors="replace"), "latin-1"


# ── Format-aware text extraction ───────────────────────────────────────────────

# Extensions handled by binary extractors — raw bytes of these must NOT be
# passed to _decode_bytes as they are ZIP/binary container formats.
_BINARY_EXTS = {
    "docx", "doc", "xlsx", "xls", "xlsm",
    "pptx", "ppt", "pdf", "odt", "ods", "odp",
    "epub", "rtf",
}

# Extensions that are plain text — decoded directly
_TEXT_EXTS = {
    "txt", "md", "rst", "csv", "tsv", "json", "jsonl",
    "xml", "html", "htm", "yaml", "yml", "toml", "ini",
    "cfg", "conf", "env", "log", "sql",
    # Code
    "py", "js", "ts", "jsx", "tsx", "css", "scss", "less",
    "html", "sh", "bash", "zsh", "fish", "ps1",
    "c", "cpp", "h", "hpp", "cs", "java", "kt", "swift",
    "go", "rs", "rb", "php", "lua", "r", "m",
    "tf", "tfvars", "hcl",
    "Dockerfile", "makefile", "cmake",
}


def _extract_text(ext: str, raw: bytes, path: str = "") -> tuple[str, str]:
    """
    Extract readable text from raw bytes.

    Returns (text, format_label).  Raises on unrecoverable extraction failure.
    Callers should pass the full raw bytes (not pre-truncated) for binary formats
    so the container parser has valid data; text truncation happens afterwards.
    """
    ext = ext.lower().lstrip(".")

    # ── Plain-text formats ────────────────────────────────────────────────────
    if ext == "json" or ext == "jsonl":
        text, enc = _decode_bytes(raw)
        try:
            import json as _json
            parsed = _json.loads(text)
            return _json.dumps(parsed, indent=2, ensure_ascii=False), "json"
        except Exception:
            return text, enc

    if ext in ("csv", "tsv"):
        text, enc = _decode_bytes(raw)
        import csv as _csv, io as _io
        delim = "\t" if ext == "tsv" else ","
        try:
            rows = list(_csv.reader(_io.StringIO(text), delimiter=delim))
            if not rows:
                return text, enc
            col_widths = [max(len(str(r[i])) for r in rows if i < len(r)) for i in range(len(rows[0]))]
            lines = []
            for ri, row in enumerate(rows):
                line = " | ".join(str(cell).ljust(col_widths[i]) for i, cell in enumerate(row))
                lines.append(line)
                if ri == 0:
                    lines.append("-+-".join("-" * w for w in col_widths))
            return "\n".join(lines), enc
        except Exception:
            return text, enc

    if ext in ("html", "htm"):
        text, enc = _decode_bytes(raw)
        try:
            from html.parser import HTMLParser
            class _Strip(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self._parts: list[str] = []
                    self._skip = False
                def handle_starttag(self, tag, attrs):
                    if tag in ("script", "style"):
                        self._skip = True
                def handle_endtag(self, tag):
                    if tag in ("script", "style"):
                        self._skip = False
                def handle_data(self, data):
                    if not self._skip:
                        stripped = data.strip()
                        if stripped:
                            self._parts.append(stripped)
            p = _Strip()
            p.feed(text)
            return "\n".join(p._parts), "html→text"
        except Exception:
            return text, enc

    if ext in ("rtf",):
        text, enc = _decode_bytes(raw)
        # Strip RTF control words with a simple regex pass
        import re as _re
        text = _re.sub(r"\\\w+\s?", "", text)
        text = _re.sub(r"[{}]", "", text)
        return text.strip(), "rtf→text"

    # ── PDF ───────────────────────────────────────────────────────────────────
    if ext == "pdf":
        try:
            from pypdf import PdfReader
            import io as _io
            reader = PdfReader(_io.BytesIO(raw))
            pages = []
            for i, page in enumerate(reader.pages, 1):
                page_text = (page.extract_text() or "").strip()
                if page_text:
                    pages.append(f"--- Page {i} ---\n{page_text}")
            return "\n\n".join(pages) or "[No extractable text found in PDF]", "pdf→text"
        except ImportError:
            raise ImportError("pypdf not installed — run: uv add pypdf")

    # ── Word / DOCX ───────────────────────────────────────────────────────────
    if ext in ("docx", "doc"):
        try:
            from docx import Document
            import io as _io
            doc = Document(_io.BytesIO(raw))
            parts: list[str] = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    # Preserve heading levels
                    if para.style.name.startswith("Heading"):
                        level = para.style.name.split()[-1] if para.style.name.split()[-1].isdigit() else "1"
                        parts.append(f"{'#' * int(level)} {text}")
                    else:
                        parts.append(text)
            # Also extract tables
            for table in doc.tables:
                rows = []
                for row in table.rows:
                    rows.append(" | ".join(c.text.strip() for c in row.cells))
                if rows:
                    parts.append("\n".join(rows))
            return "\n\n".join(parts) or "[No extractable text found in document]", "docx→text"
        except ImportError:
            raise ImportError("python-docx not installed — run: uv add python-docx")

    # ── Excel / XLSX ──────────────────────────────────────────────────────────
    if ext in ("xlsx", "xls", "xlsm"):
        try:
            import openpyxl, io as _io
            wb = openpyxl.load_workbook(_io.BytesIO(raw), read_only=True, data_only=True)
            sheets: list[str] = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join("" if v is None else str(v) for v in row)
                    if row_text.strip():
                        rows.append(row_text)
                if rows:
                    sheets.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows))
            return "\n\n".join(sheets) or "[No extractable data found in workbook]", "xlsx→text"
        except ImportError:
            # Fall back to xlrd for .xls
            try:
                import xlrd, io as _io
                wb = xlrd.open_workbook(file_contents=raw)
                sheets = []
                for i in range(wb.nsheets):
                    ws = wb.sheet_by_index(i)
                    rows = [" | ".join(str(ws.cell_value(r, c)) for c in range(ws.ncols))
                            for r in range(ws.nrows)]
                    sheets.append(f"=== Sheet: {ws.name} ===\n" + "\n".join(rows))
                return "\n\n".join(sheets), "xls→text"
            except ImportError:
                raise ImportError("openpyxl not installed — run: uv add openpyxl")

    # ── PowerPoint / PPTX ────────────────────────────────────────────────────
    if ext in ("pptx", "ppt"):
        try:
            from pptx import Presentation
            import io as _io
            prs = Presentation(_io.BytesIO(raw))
            slides: list[str] = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                if texts:
                    slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))
            return "\n\n".join(slides) or "[No extractable text found in presentation]", "pptx→text"
        except ImportError:
            raise ImportError("python-pptx not installed — run: uv add python-pptx")

    # ── EPUB ─────────────────────────────────────────────────────────────────
    if ext == "epub":
        try:
            import zipfile, io as _io
            from html.parser import HTMLParser
            class _Strip(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self._parts: list[str] = []
                    self._skip = False
                def handle_starttag(self, tag, attrs):
                    if tag in ("script", "style"): self._skip = True
                def handle_endtag(self, tag):
                    if tag in ("script", "style"): self._skip = False
                def handle_data(self, data):
                    if not self._skip:
                        s = data.strip()
                        if s: self._parts.append(s)
            with zipfile.ZipFile(_io.BytesIO(raw)) as zf:
                parts: list[str] = []
                for name in sorted(zf.namelist()):
                    if name.endswith((".html", ".xhtml", ".htm")):
                        html_bytes = zf.read(name)
                        html_text = html_bytes.decode("utf-8", errors="replace")
                        p = _Strip()
                        p.feed(html_text)
                        parts.extend(p._parts)
            return "\n\n".join(parts) or "[No extractable text found in EPUB]", "epub→text"
        except Exception as e:
            raise ValueError(f"Could not read EPUB: {e}")

    # ── Default: treat as plain text ──────────────────────────────────────────
    return _decode_bytes(raw)


# ── Human-readable helpers ──────────────────────────────────────────────────────

def _human_size(b: int) -> str:
    for unit in ("B", "KB", "MB", "GB", "TB"):
        if b < 1024:
            return f"{b:.1f} {unit}" if unit != "B" else f"{b} B"
        b //= 1024
    return f"{b:.1f} TB"


def _iso(ts: float) -> str:
    import datetime
    return datetime.datetime.fromtimestamp(ts).isoformat(timespec="seconds")


def _perm_human(mode: int) -> str:
    """Convert stat mode to rwxrwxrwx string."""
    chars = ""
    for shift in (6, 3, 0):
        bits = (mode >> shift) & 0o7
        chars += "r" if bits & 4 else "-"
        chars += "w" if bits & 2 else "-"
        chars += "x" if bits & 1 else "-"
    return chars


def _is_hidden(path: Path) -> bool:
    """Cross-platform hidden-file detection."""
    if path.name.startswith("."):
        return True
    if os.name == "nt":
        try:
            attrs = path.stat().st_file_attributes  # type: ignore[attr-defined]
            return bool(attrs & stat.FILE_ATTRIBUTE_HIDDEN)  # type: ignore[attr-defined]
        except Exception:
            pass
    return False


# ── OS Documents folder ────────────────────────────────────────────────────────

def _get_documents_dir() -> Path:
    """Return the user's Documents folder cross-platform."""
    if os.name == "nt":
        try:
            import ctypes, ctypes.wintypes
            buf = ctypes.create_unicode_buffer(ctypes.wintypes.MAX_PATH)
            # CSIDL_PERSONAL (0x05) = My Documents
            ctypes.windll.shell32.SHGetFolderPathW(0, 0x05, 0, 0, buf)  # type: ignore[attr-defined]
            docs = Path(buf.value)
            if docs.exists():
                return docs
        except Exception:
            pass
    # macOS / Linux / fallback
    docs = Path.home() / "Documents"
    docs.mkdir(parents=True, exist_ok=True)
    return docs


# ── Max sizes ───────────────────────────────────────────────────────────────────

_MAX_READ_DEFAULT_MB = 1
_MAX_READ_HARD_MB    = 10


# ── File operations ─────────────────────────────────────────────────────────────

class FileOps:
    """All filesystem operations used by the File R/W skill."""

    def __init__(self, config: dict | None = None):
        cfg = config or {}
        self.max_read_bytes    = int(cfg.get("max_read_size_mb", _MAX_READ_DEFAULT_MB)) * 1024 * 1024
        self.create_backups    = bool(cfg.get("create_backups", True))
        self.default_permanent = cfg.get("default_delete_mode", "recycle_bin") == "permanent"

        # ── Default create path ────────────────────────────────────────────────
        configured = cfg.get("default_create_path", "").strip()
        if configured:
            self.default_create_path = Path(configured)
        else:
            self.default_create_path = _get_documents_dir()

    # ── End __init__ ────────────────────────────────────────────────────────────

    # ── READ ────────────────────────────────────────────────────────────────────

    async def read_file(self, path: str) -> FileReadResult:
        return await asyncio.get_event_loop().run_in_executor(
            None, self._read_sync, path
        )

    def _read_sync(self, path: str) -> FileReadResult:
        p = Path(path)
        if not p.exists():
            raise FileNotFoundError(f"File not found: {path}")
        if not p.is_file():
            raise IsADirectoryError(f"Path is a directory: {path}")

        size = p.stat().st_size
        hard_limit = _MAX_READ_HARD_MB * 1024 * 1024

        if size > hard_limit:
            raise ValueError(
                f"File too large to read ({_human_size(size)} > {_MAX_READ_HARD_MB} MB). "
                "Use the stats operation to inspect it instead."
            )

        raw = p.read_bytes()
        was_truncated = False
        truncated_at = 0
        ext = p.suffix.lstrip(".").lower()

        if ext in _BINARY_EXTS:
            # Binary containers: parse the FULL bytes with format-aware extractor,
            # then text-truncate the *extracted* result (not the raw bytes).
            content, encoding = _extract_text(ext, raw, path)
            if len(content) > self.max_read_bytes:
                content = content[: self.max_read_bytes]
                was_truncated = True
                truncated_at = self.max_read_bytes
                content += (
                    f"\n\n[Text truncated at {_human_size(self.max_read_bytes)} — "
                    f"source file: {_human_size(size)}]"
                )
        else:
            # Plain-text formats: byte-truncate first, then decode / format-extract.
            if size > self.max_read_bytes:
                raw = raw[: self.max_read_bytes]
                was_truncated = True
                truncated_at = self.max_read_bytes

            content, encoding = _extract_text(ext, raw, path)
            if was_truncated:
                content += (
                    f"\n\n[File truncated at {_human_size(self.max_read_bytes)} — "
                    f"total size: {_human_size(size)}]"
                )

        return FileReadResult(
            content=content,
            encoding=encoding,
            size_bytes=size,
            line_count=content.count("\n") + 1,
            was_truncated=was_truncated,
            truncated_at=truncated_at,
        )

    # ── WRITE (full overwrite) ──────────────────────────────────────────────────

    async def read_original_for_diff(self, path: str) -> str:
        """Read existing content (empty string if file does not exist)."""
        try:
            result = await self.read_file(path)
            return result.content
        except (FileNotFoundError, IsADirectoryError):
            return ""

    async def write_file(
        self, path: str, content: str
    ) -> WriteResult:
        """Atomically write content to path, creating a backup first."""
        return await asyncio.get_event_loop().run_in_executor(
            None, self._write_sync, path, content
        )

    def _write_sync(self, path: str, content: str) -> WriteResult:
        p = Path(path)
        backup_path: Optional[str] = None

        # Backup
        if self.create_backups and p.exists():
            bk = Path(str(path) + ".ragdoll_backup")
            shutil.copy2(path, bk)
            backup_path = str(bk)

        # Atomic write: write to .tmp, then rename
        tmp = Path(str(path) + ".tmp")
        try:
            p.parent.mkdir(parents=True, exist_ok=True)
            tmp.write_text(content, encoding="utf-8")
            os.replace(tmp, path)
        except Exception as exc:
            tmp.unlink(missing_ok=True)
            return WriteResult(ok=False, path=path, error=str(exc))

        return WriteResult(ok=True, path=path, backup_path=backup_path)

    # ── PATCH (targeted edits) ──────────────────────────────────────────────────

    def apply_patches(
        self, content: str, patches: list[dict]
    ) -> tuple[str, list[str]]:
        """
        Apply a list of {search, replace, occurrence} patches to *content*.
        Returns (patched_content, list_of_errors).
        occurrence: 1=first, -1=last, 0=all
        """
        errors: list[str] = []
        result = content
        for patch in patches:
            search  = patch.get("search", "")
            replace = patch.get("replace", "")
            occ     = int(patch.get("occurrence", 1))

            if not search:
                errors.append("Patch has empty 'search' field — skipped.")
                continue
            if search not in result:
                errors.append(
                    f"Search string not found: {search[:60]!r} — skipped."
                )
                continue

            if occ == 0:
                result = result.replace(search, replace)
            elif occ == -1:
                idx = result.rfind(search)
                result = result[:idx] + replace + result[idx + len(search):]
            else:
                # Replace nth occurrence (1-based)
                count = 0
                pos = 0
                while True:
                    idx = result.find(search, pos)
                    if idx == -1:
                        errors.append(
                            f"Occurrence {occ} of {search[:40]!r} not found."
                        )
                        break
                    count += 1
                    if count == occ:
                        result = result[:idx] + replace + result[idx + len(search):]
                        break
                    pos = idx + len(search)

        return result, errors

    # ── CREATE ──────────────────────────────────────────────────────────────────

    async def create_file(
        self, path: str, content: str = "", overwrite: bool = False
    ) -> OperationResult:
        p = Path(path)
        if not p.is_absolute():
            p = self.default_create_path / p
        path = str(p)
        if p.exists() and not overwrite:
            return OperationResult(
                ok=False, path=path,
                error=f"File already exists: {path}. Use overwrite=true to replace it."
            )
        try:
            p.parent.mkdir(parents=True, exist_ok=True)
            ext = p.suffix.lstrip(".").lower()
            from doc_formatter import FORMATTED_EXTS, render_document
            if ext in FORMATTED_EXTS:
                # Rich format — Markdown content is rendered into the native format
                render_document(ext, content, path)
            else:
                # Plain text — atomic write
                tmp = Path(str(path) + ".tmp")
                tmp.write_text(content, encoding="utf-8")
                os.replace(tmp, path)
            return OperationResult(ok=True, path=path)
        except Exception as exc:
            return OperationResult(ok=False, path=path, error=str(exc))

    async def create_directory(self, path: str) -> OperationResult:
        p = Path(path)
        if not p.is_absolute():
            p = self.default_create_path / p
        path = str(p)
        try:
            p.mkdir(parents=True, exist_ok=True)
            return OperationResult(ok=True, path=path)
        except Exception as exc:
            return OperationResult(ok=False, path=path, error=str(exc))

    # ── DELETE ──────────────────────────────────────────────────────────────────

    async def delete(
        self, path: str, permanent: bool = False
    ) -> OperationResult:
        return await asyncio.get_event_loop().run_in_executor(
            None, self._delete_sync, path, permanent
        )

    def _delete_sync(self, path: str, permanent: bool) -> OperationResult:
        p = Path(path)
        if not p.exists():
            return OperationResult(
                ok=False, path=path, error=f"Path does not exist: {path}"
            )
        try:
            if permanent:
                if p.is_dir():
                    shutil.rmtree(path)
                else:
                    p.unlink()
                mode = "permanently deleted"
            else:
                import send2trash
                send2trash.send2trash(str(p.resolve()))
                mode = "moved to trash"
            return OperationResult(ok=True, path=path, extra={"mode": mode})
        except ImportError:
            # send2trash not available — fall back to permanent
            if p.is_dir():
                shutil.rmtree(path)
            else:
                p.unlink()
            return OperationResult(
                ok=True, path=path,
                extra={"mode": "permanently deleted (send2trash unavailable)"},
            )
        except Exception as exc:
            return OperationResult(ok=False, path=path, error=str(exc))

    async def delete_multiple(
        self, paths: list[str], permanent: bool = False
    ) -> list[OperationResult]:
        # Sequential — safer for filesystem
        results = []
        for p in paths:
            results.append(await self.delete(p, permanent))
        return results

    # ── RENAME / MOVE / COPY ────────────────────────────────────────────────────

    async def rename(self, old_path: str, new_name: str) -> OperationResult:
        try:
            old = Path(old_path)
            new = old.parent / new_name
            if new.exists():
                return OperationResult(
                    ok=False, path=old_path,
                    error=f"Destination already exists: {new}"
                )
            os.rename(old, new)
            return OperationResult(ok=True, path=str(new))
        except Exception as exc:
            return OperationResult(ok=False, path=old_path, error=str(exc))

    async def move(self, src: str, dst: str) -> OperationResult:
        try:
            shutil.move(src, dst)
            return OperationResult(ok=True, path=dst)
        except Exception as exc:
            return OperationResult(ok=False, path=src, error=str(exc))

    async def copy(self, src: str, dst: str) -> OperationResult:
        try:
            p = Path(src)
            if p.is_dir():
                shutil.copytree(src, dst)
            else:
                shutil.copy2(src, dst)
            return OperationResult(ok=True, path=dst)
        except Exception as exc:
            return OperationResult(ok=False, path=src, error=str(exc))

    # ── STATS ────────────────────────────────────────────────────────────────────

    async def get_stats(self, path: str) -> FileStats:
        return await asyncio.get_event_loop().run_in_executor(
            None, self._stats_sync, path
        )

    def _stats_sync(self, path: str) -> FileStats:
        p = Path(path).resolve()
        if not p.exists():
            raise FileNotFoundError(f"Path does not exist: {path}")

        s = p.stat()
        mime, _ = mimetypes.guess_type(str(p))

        # Owner
        try:
            import pwd
            owner = pwd.getpwuid(s.st_uid).pw_name
        except Exception:
            try:
                import subprocess
                owner = subprocess.check_output(
                    ["icacls", str(p)], text=True, stderr=subprocess.DEVNULL
                ).splitlines()[0].split()[-1] if os.name == "nt" else str(s.st_uid)
            except Exception:
                owner = str(getattr(s, "st_uid", "unknown"))

        child_count = total_size = None
        if p.is_dir():
            child_count = 0
            total_size  = 0
            deadline = time.time() + 10
            try:
                for entry in os.scandir(p):
                    child_count += 1
                    if time.time() > deadline:
                        break
                    try:
                        es = entry.stat()
                        total_size += es.st_size
                    except OSError:
                        pass
            except PermissionError:
                pass

        return FileStats(
            name=p.name,
            absolute_path=str(p),
            size_bytes=s.st_size,
            size_human=_human_size(s.st_size),
            is_file=p.is_file(),
            is_directory=p.is_dir(),
            is_symlink=p.is_symlink(),
            is_hidden=_is_hidden(p),
            extension=p.suffix.lstrip("."),
            mime_type=mime or "",
            created_at=_iso(getattr(s, "st_birthtime", s.st_ctime)),
            modified_at=_iso(s.st_mtime),
            accessed_at=_iso(s.st_atime),
            permissions_octal=oct(stat.S_IMODE(s.st_mode)),
            permissions_human=_perm_human(stat.S_IMODE(s.st_mode)),
            owner=owner,
            is_readable=os.access(p, os.R_OK),
            is_writable=os.access(p, os.W_OK),
            is_executable=os.access(p, os.X_OK),
            child_count=child_count,
            total_size=total_size,
        )

    # ── LIST ─────────────────────────────────────────────────────────────────────

    async def list_directory(
        self,
        path: str,
        show_hidden: bool = False,
        sort_by: str = "name",
    ) -> DirectoryListing:
        return await asyncio.get_event_loop().run_in_executor(
            None, self._list_sync, path, show_hidden, sort_by
        )

    def _list_sync(
        self, path: str, show_hidden: bool, sort_by: str
    ) -> DirectoryListing:
        p = Path(path)
        if not p.is_dir():
            raise NotADirectoryError(f"Not a directory: {path}")

        entries: list[DirectoryEntry] = []
        hidden_count = 0

        for entry in os.scandir(p):
            ep = Path(entry.path)
            hidden = _is_hidden(ep)
            if hidden:
                hidden_count += 1
            if hidden and not show_hidden:
                continue
            try:
                es = entry.stat()
                size = es.st_size if entry.is_file() else None
                mod  = _iso(es.st_mtime)
            except OSError:
                size = None
                mod  = ""
            entries.append(DirectoryEntry(
                name=entry.name,
                is_directory=entry.is_dir(),
                size_bytes=size,
                modified_at=mod,
                extension=ep.suffix.lstrip("."),
                is_hidden=hidden,
            ))

        # Sort
        key_fns = {
            "name":     lambda e: (e.is_directory, e.name.lower()),
            "modified": lambda e: e.modified_at,
            "size":     lambda e: e.size_bytes or 0,
            "type":     lambda e: (e.extension.lower(), e.name.lower()),
        }
        entries.sort(key=key_fns.get(sort_by, key_fns["name"]))

        return DirectoryListing(
            path=str(p),
            entries=entries,
            total_count=len(entries) + hidden_count,
            hidden_count=hidden_count,
        )

    # ── UNDO ─────────────────────────────────────────────────────────────────────

    async def undo_write(self, path: str) -> tuple[bool, str]:
        """
        Restore the .ragdoll_backup file to *path*.
        Returns (success, message).
        """
        bk = Path(str(path) + ".ragdoll_backup")
        if not bk.exists():
            return False, "No backup available for this file."
        try:
            tmp = Path(str(path) + ".undo_tmp")
            shutil.copy2(bk, tmp)
            os.replace(tmp, path)
            return True, f"Restored backup to {path}"
        except Exception as exc:
            return False, str(exc)
